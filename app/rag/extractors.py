"""多格式文档解析：把各类文件统一提取为纯文本，交给现有切分/向量管线。

支持：txt / md / pdf / docx / xlsx / csv / pptx。
解析失败会抛出异常，由调用方（indexer）决定跳过该文件。
"""
from __future__ import annotations

import csv
import os
import json
import re
import zipfile
import xml.etree.ElementTree as ET
from html.parser import HTMLParser

SUPPORTED_EXTS = {".txt", ".md", ".pdf", ".docx", ".xlsx", ".csv", ".pptx",
                  ".html", ".json", ".rtf", ".xml", ".xls", ".odt", ".epub"}


def _rows_to_markdown(header: list, rows: list) -> str:
    """把二维行数据渲染成 Markdown 表格，供结构感知切分识别、表头保持可见。"""

    def esc(cell) -> str:
        return str(cell if cell is not None else "").replace("|", "\\|").replace("\n", " ")

    cols = [esc(c) for c in header]
    lines = ["| " + " | ".join(cols) + " |", "|" + "---|" * len(cols)]
    for row in rows:
        cells = [esc(c) for c in row]
        while len(cells) < len(cols):
            cells.append("")
        lines.append("| " + " | ".join(cells) + " |")
    return "\n".join(lines)


def _field_list_line(header: list) -> str:
    """输出「字段列表：...」一行，让字段名直接可检索。"""
    names = [str(c).strip() for c in header if str(c).strip()]
    return "字段列表：" + "、".join(names) if names else ""


def _read_txt(path: str) -> str:
    for enc in ("utf-8-sig", "utf-8", "gbk"):
        try:
            with open(path, "r", encoding=enc) as f:
                return f.read()
        except UnicodeDecodeError:
            continue
    with open(path, "r", encoding="utf-8", errors="ignore") as f:
        return f.read()


def _read_pdf(path: str) -> str:
    from pypdf import PdfReader

    reader = PdfReader(path)
    parts = []
    for i, page in enumerate(reader.pages, 1):
        text = page.extract_text() or ""
        parts.append(f"[第{i}页]\n{text}")
    return "\n\n".join(parts)


def _read_docx(path: str) -> str:
    from docx import Document

    doc = Document(path)
    parts = [p.text for p in doc.paragraphs if p.text.strip()]
    for table in doc.tables:
        rows = [[cell.text.strip() for cell in row.cells] for row in table.rows]
        rows = [r for r in rows if any(r)]
        if rows:
            parts.append(_rows_to_markdown(rows[0], rows[1:]))
    return "\n".join(parts)


def _read_xlsx(path: str) -> str:
    from openpyxl import load_workbook

    wb = load_workbook(path, read_only=True, data_only=True)
    parts = []
    for ws in wb.worksheets:
        rows = []
        for row in ws.iter_rows(values_only=True):
            vals = ["" if v is None else str(v) for v in row]
            if any(vals):
                rows.append(vals)
        if not rows:
            continue
        header = rows[0]
        block = [f"# Sheet: {ws.title}", _field_list_line(header), _rows_to_markdown(header, rows[1:])]
        parts.append("\n\n".join(p for p in block if p))
    return "\n\n".join(parts)


def _read_csv(path: str) -> str:
    with open(path, "r", encoding="utf-8-sig", newline="") as f:
        rows = [row for row in csv.reader(f) if any(cell.strip() for cell in row)]
    if not rows:
        return ""
    header = rows[0]
    return "\n\n".join(
        p for p in (_field_list_line(header), _rows_to_markdown(header, rows[1:])) if p
    )


def _read_pptx(path: str) -> str:
    from pptx import Presentation

    prs = Presentation(path)
    parts = []
    for i, slide in enumerate(prs.slides, 1):
        slide_lines = [f"[第{i}页]"]
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = "\n".join(
                    p.text for p in shape.text_frame.paragraphs if p.text.strip()
                )
                if text:
                    slide_lines.append(text)
            if getattr(shape, "has_table", False):
                rows = [[cell.text.strip() for cell in row.cells] for row in shape.table.rows]
                rows = [r for r in rows if any(r)]
                if rows:
                    slide_lines.append(_rows_to_markdown(rows[0], rows[1:]))
        parts.append("\n".join(slide_lines))
    return "\n\n".join(parts)


class _HTMLTextExtractor(HTMLParser):
    """剥掉 HTML 标签/脚本，保留正文与表格分隔。"""

    def __init__(self) -> None:
        super().__init__()
        self.parts: list[str] = []
        self.skip = 0

    def handle_starttag(self, tag, attrs) -> None:
        if tag in ("script", "style"):
            self.skip += 1
        if tag in ("p", "div", "tr", "h1", "h2", "h3", "h4", "li", "br",
                   "table", "section", "article"):
            self.parts.append("\n")
        if tag in ("td", "th"):
            self.parts.append(" | ")

    def handle_endtag(self, tag) -> None:
        if tag in ("script", "style") and self.skip:
            self.skip -= 1

    def handle_data(self, data) -> None:
        if not self.skip:
            self.parts.append(data)


def _html_to_text(raw: str) -> str:
    p = _HTMLTextExtractor()
    try:
        p.feed(raw)
    except Exception:
        pass
    text = "".join(p.parts)
    lines = [ln.strip() for ln in text.splitlines()]
    return "\n".join(ln for ln in lines if ln)


def _read_html(path: str) -> str:
    with open(path, "r", encoding="utf-8", errors="ignore") as f:
        return _html_to_text(f.read())


# 键值/配置类文件的「中文语义标签」：让「端口」能匹配 port、「最大连接数」能匹配 max_connections。
# 检索是词面+向量匹配，原始英文键名与中文问法不在同一语言空间，必须在这里补标签。
_JSON_KEY_LABELS = {
    "port": "端口",
    "host": "主机",
    "url": "地址",
    "endpoint": "接口地址",
    "timeout": "超时时间",
    "maxconnections": "最大连接数(并发连接)",
    "connection": "连接数",
    "connections": "连接数",
    "qps": "每秒请求数(QPS)",
    "ratelimit": "限流",
    "version": "版本",
    "service": "服务",
    "name": "名称",
    "database": "数据库",
    "redis": "缓存",
    "cache": "缓存",
    "api": "接口",
    "price": "价格",
    "cost": "费用",
    "stock": "库存",
    "quantity": "数量",
    "storage": "存储",
    "bucket": "存储桶",
    "model": "模型",
    "key": "密钥",
    "secret": "密钥",
    "token": "令牌",
    "user": "用户",
    "password": "密码",
    "username": "用户名",
    "email": "邮箱",
    "phone": "电话",
    "address": "地址",
    "region": "区域",
    "status": "状态",
    "enabled": "启用",
    "disabled": "禁用",
    "featureflags": "功能开关",
    "multimodal": "多模态",
    "mcpmarket": "MCP工具市场",
    "topk": "检索条数(top_k)",
    "rerank": "重排",
    "embeddingmodel": "向量模型",
    "type": "类型",
    "count": "数量",
    "total": "总数",
    "limit": "上限",
    "max": "最大值",
    "min": "最小值",
    "interval": "间隔",
    "ttl": "有效期",
}

_JSON_FULL_KEY_LABELS = {
    "database.max_connections": "数据库最大连接数(并发连接)",
    "database.host": "数据库主机",
    "database.port": "数据库端口",
    "database.type": "数据库类型",
    "database.name": "数据库名称",
    "rate_limit.qps": "限流每秒请求数(QPS)",
    "rate_limit.burst": "限流突发上限",
    "rag.top_k": "检索返回条数(top_k)",
    "rag.rerank": "检索重排开关",
    "rag.embedding_model": "向量模型",
    "feature_flags.multimodal": "多模态开关",
    "feature_flags.mcp_market": "MCP工具市场开关",
    "service.port": "服务端口",
    "service.version": "服务版本",
}


def _json_key_label(path: tuple) -> str:
    """按完整键路径优先、叶子键兜底，返回中文标签；没有则返回空串。"""
    full = ".".join(path)
    if full in _JSON_FULL_KEY_LABELS:
        return _JSON_FULL_KEY_LABELS[full]
    leaf = path[-1]
    norm = re.sub(r"[^a-z0-9]", "", leaf.lower())
    return _JSON_KEY_LABELS.get(norm, "")


def _flatten_json(obj, prefix: tuple = (), out: list | None = None) -> list:
    out = out if out is not None else []
    if isinstance(obj, dict):
        for k, v in obj.items():
            _flatten_json(v, prefix + (k,), out)
    elif isinstance(obj, list):
        for i, v in enumerate(obj):
            _flatten_json(v, prefix + (str(i),), out)
    else:
        key = ".".join(prefix)
        label = _json_key_label(prefix) if prefix else ""
        out.append(f"{label}({key}): {obj}" if label else f"{key}: {obj}")
    return out


def _read_json(path: str) -> str:
    with open(path, "r", encoding="utf-8") as f:
        data = json.load(f)
    return "\n".join(_flatten_json(data))


def _read_rtf(path: str) -> str:
    with open(path, "r", encoding="utf-8", errors="ignore") as f:
        raw = f.read()
    raw = re.sub(r"\{\\fonttbl[^}]*\}", "", raw, flags=re.I)  # 去掉字体表
    # \uN? -> 字符
    raw = re.sub(r"\\u(-?\d+)\?", lambda m: chr(max(0, int(m.group(1)))), raw)
    raw = re.sub(r"\\par\b|\\line\b|\\\n", "\n", raw, flags=re.I)
    raw = re.sub(r"\\'[0-9a-fA-F]{2}", "", raw)
    raw = re.sub(r"\\[a-zA-Z]+-?\d* ?", "", raw)
    raw = raw.replace("{", "").replace("}", "")
    lines = [ln.strip() for ln in raw.splitlines()]
    return "\n".join(ln for ln in lines if ln)


def _read_xml(path: str) -> str:
    tree = ET.parse(path)
    lines = []
    for elem in tree.iter():
        if elem.attrib:
            lines.append(f"<{elem.tag}> " + " ".join(
                f"{k}={v}" for k, v in elem.attrib.items()))
        if elem.text and elem.text.strip():
            lines.append(f"{elem.tag}: {elem.text.strip()}")
    return "\n".join(lines)


def _read_xls(path: str) -> str:
    import xlrd

    book = xlrd.open_workbook(path)
    parts = []
    for sh in book.sheets():
        rows = []
        for r in range(sh.nrows):
            vals = ["" if sh.cell_value(r, c) is None else str(sh.cell_value(r, c))
                    for c in range(sh.ncols)]
            if any(vals):
                rows.append(vals)
        if not rows:
            continue
        header = rows[0]
        block = [f"# Sheet: {sh.name}", _field_list_line(header), _rows_to_markdown(header, rows[1:])]
        parts.append("\n\n".join(p for p in block if p))
    return "\n\n".join(parts)


def _read_odt(path: str) -> str:
    with zipfile.ZipFile(path) as z:
        root = ET.fromstring(z.read("content.xml"))
    parts = [e.text.strip() for e in root.iter() if e.text and e.text.strip()]
    return "\n".join(parts)


def _read_epub(path: str) -> str:
    with zipfile.ZipFile(path) as z:
        names = [n for n in z.namelist()
                 if n.lower().endswith((".xhtml", ".html", ".htm"))]
        parts = []
        for n in sorted(names):
            raw = z.read(n).decode("utf-8", errors="ignore")
            text = _html_to_text(raw)
            if text.strip():
                parts.append(f"[{n}]\n{text}")
    return "\n\n".join(parts)


_EXTRACTORS = {
    ".txt": _read_txt,
    ".md": _read_txt,
    ".pdf": _read_pdf,
    ".docx": _read_docx,
    ".xlsx": _read_xlsx,
    ".csv": _read_csv,
    ".pptx": _read_pptx,
    ".html": _read_html,
    ".json": _read_json,
    ".rtf": _read_rtf,
    ".xml": _read_xml,
    ".xls": _read_xls,
    ".odt": _read_odt,
    ".epub": _read_epub,
}


def extract_text(path: str) -> str:
    """按扩展名提取文本；不支持的格式抛 ValueError。"""
    ext = os.path.splitext(path)[1].lower()
    fn = _EXTRACTORS.get(ext)
    if fn is None:
        raise ValueError(f"unsupported file type: {ext}")
    return fn(path)
